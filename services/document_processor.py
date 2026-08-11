"""
Campus Cognition V2 — Document Processing Service
Handles extraction of text from PDF, DOCX, PPTX, and TXT files.
Includes hashing for deduplication and caching.
"""
import os
import hashlib
import logging
import mimetypes
import PyPDF2
from typing import Optional, Tuple
import docx
from pptx import Presentation

logger = logging.getLogger(__name__)

def generate_file_hash(filepath: str) -> str:
    """Generate SHA-256 hash of a file for caching purposes."""
    hasher = hashlib.sha256()
    with open(filepath, 'rb') as f:
        buf = f.read(65536)
        while len(buf) > 0:
            hasher.update(buf)
            buf = f.read(65536)
    return hasher.hexdigest()

def extract_text_from_pdf(filepath: str, max_chars: int = 15000) -> str:
    """Extract text from a PDF file up to max_chars."""
    text = ""
    try:
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                if len(text) >= max_chars:
                    break
    except Exception as e:
        logger.error(f"Error reading PDF {filepath}: {e}")
    return text[:max_chars]

def extract_text_from_docx(filepath: str, max_chars: int = 15000) -> str:
    """Extract text from a DOCX file."""
    text = ""
    try:
        doc = docx.Document(filepath)
        for para in doc.paragraphs:
            text += para.text + "\n"
            if len(text) >= max_chars:
                break
    except Exception as e:
        logger.error(f"Error reading DOCX {filepath}: {e}")
    return text[:max_chars]

def extract_text_from_pptx(filepath: str, max_chars: int = 15000) -> str:
    """Extract text from a PPTX file."""
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            if len(text) >= max_chars:
                break
    except Exception as e:
        logger.error(f"Error reading PPTX {filepath}: {e}")
    return text[:max_chars]

def extract_text_from_txt(filepath: str, max_chars: int = 15000) -> str:
    """Extract text from a TXT/MD file."""
    text = ""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read(max_chars)
    except Exception as e:
        logger.error(f"Error reading TXT {filepath}: {e}")
    return text

def validate_document_file(filepath: str, max_size_mb: int = 15) -> str:
    """
    Validates a file based on size and magic bytes (file signature), not just extension.
    Returns the resolved extension ('pdf', 'docx', 'pptx', 'txt', 'md').
    Raises ValueError if invalid or corrupted.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found: {filepath}")
        
    file_size = os.path.getsize(filepath)
    if file_size > max_size_mb * 1024 * 1024:
        raise ValueError(f"File exceeds maximum allowed size of {max_size_mb}MB")
        
    # Read first 8 bytes for signature
    with open(filepath, 'rb') as f:
        header = f.read(8)
        
    ext = filepath.rsplit('.', 1)[-1].lower() if '.' in filepath else ''
    
    # 1. PDF Signature: %PDF-
    if header.startswith(b'%PDF-'):
        return 'pdf'
        
    # 2. ZIP Signature (DOCX, PPTX): PK\x03\x04
    if header.startswith(b'PK\x03\x04'):
        if ext == 'docx':
            return 'docx'
        elif ext == 'pptx':
            return 'pptx'
        else:
            raise ValueError("Invalid ZIP-based document type")
            
    # 3. Fallback to text checks for TXT/MD
    if ext in ['txt', 'md']:
        # Basic check to see if it's text (no null bytes in first block)
        if b'\x00' in header:
            raise ValueError("Corrupted or invalid text file")
        return ext
        
    raise ValueError("Unsupported or corrupted file format")


def process_document(filepath: str, max_chars: int = 15000) -> Tuple[str, str]:
    """
    Process a document and return (file_hash, extracted_text).
    Supports .pdf, .docx, .pptx, .txt, .md with strict validation.
    """
    ext = validate_document_file(filepath)
    file_hash = generate_file_hash(filepath)
    
    if ext == 'pdf':
        text = extract_text_from_pdf(filepath, max_chars)
    elif ext == 'docx':
        text = extract_text_from_docx(filepath, max_chars)
    elif ext == 'pptx':
        text = extract_text_from_pptx(filepath, max_chars)
    elif ext in ['txt', 'md']:
        text = extract_text_from_txt(filepath, max_chars)
    else:
        # Fallback to text reading if possible
        text = extract_text_from_txt(filepath, max_chars)
        
    return file_hash, text
